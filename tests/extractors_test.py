from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING

import pytest
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.placeholder import NotesSlidePlaceholder
from pptx.slide import Slide
from pptx.text.text import TextFrame

from kreuzberg._extractors import (
    convert_pdf_to_images,
    extract_content_with_pandoc,
    extract_file_with_pandoc,
    extract_pdf_file,
    extract_pdf_with_pdfium2,
    extract_pdf_with_tesseract,
    extract_pptx_file,
)
from kreuzberg._tesseract import process_image_with_tesseract
from kreuzberg.exceptions import OCRError, ParsingError

if TYPE_CHECKING:
    from pytest_mock import MockerFixture


async def test_extract_pdf_with_pdfium2(searchable_pdf: Path) -> None:
    result = await extract_pdf_with_pdfium2(searchable_pdf)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_pdf_with_tesseract(scanned_pdf: Path) -> None:
    result = await extract_pdf_with_tesseract(scanned_pdf)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_pdf_file(searchable_pdf: Path) -> None:
    result = await extract_pdf_file(searchable_pdf)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_pdf_file_non_searchable(non_searchable_pdf: Path) -> None:
    result = await extract_pdf_file(non_searchable_pdf)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_pdf_file_invalid() -> None:
    with pytest.raises(FileNotFoundError):
        await extract_pdf_file(Path("/invalid/path.pdf"))


async def test_extract_content_with_pandoc(docx_document: Path) -> None:
    content = docx_document.read_bytes()
    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = await extract_content_with_pandoc(content, mime_type)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_file_with_pandoc(docx_document: Path) -> None:
    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    result = await extract_file_with_pandoc(docx_document, mime_type)
    assert isinstance(result, str)
    assert result.strip()


async def test_extract_file_with_pandoc_invalid() -> None:
    with pytest.raises(ParsingError):
        await extract_file_with_pandoc(
            "/invalid/path.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )


async def test_process_image_with_tesseract(ocr_image: Path) -> None:
    result = await process_image_with_tesseract(ocr_image)
    assert isinstance(result, str)
    assert result.strip()


async def test_process_image_with_tesseract_invalid() -> None:
    with pytest.raises(OCRError):
        await process_image_with_tesseract("/invalid/path.jpg")


async def test_convert_pdf_to_images_raises_parsing_error(tmp_path: Path) -> None:
    pdf_path = tmp_path / "test.pdf"
    pdf_path.write_text("invalid pdf content")

    with pytest.raises(ParsingError) as exc_info:
        await convert_pdf_to_images(pdf_path)

    assert "Could not convert PDF to images" in str(exc_info.value)
    assert str(pdf_path) in str(exc_info.value.context["file_path"])


async def test_extract_content_with_pandoc_raises_parsing_error() -> None:
    mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    with pytest.raises(ParsingError):
        await extract_content_with_pandoc(b"invalid content", mime_type)


async def test_extract_pptx_with_notes(mocker: MockerFixture) -> None:
    mock_presentation = mocker.Mock(spec=Presentation)
    mock_slide = mocker.Mock(spec=Slide)
    mock_notes_slide = mocker.Mock(spec=NotesSlidePlaceholder)
    mock_notes_frame = mocker.Mock(spec=TextFrame)
    mock_shape = mocker.Mock(spec=Shape)
    mock_shapes = mocker.Mock()

    mock_presentation.slides = [mock_slide]
    mock_slide.shapes = mock_shapes
    mock_shapes.title = None
    mock_shapes.__iter__ = lambda _: iter([mock_shape])
    mock_slide.has_notes_slide = True
    mock_slide.notes_slide = mock_notes_slide
    mock_notes_slide.notes_text_frame = mock_notes_frame
    mock_notes_frame.text = "Test note content"
    mock_shape.has_text_frame = False

    mocker.patch("pptx.Presentation", return_value=mock_presentation)

    result = await extract_pptx_file(b"mock pptx content")

    assert "Test note content" in result


async def test_extract_pdf_with_pdfium2_raises_parsing_error(tmp_path: Path) -> None:
    pdf_path = tmp_path / "invalid.pdf"
    pdf_path.write_text("invalid pdf content")

    with pytest.raises(ParsingError) as exc_info:
        await extract_pdf_with_pdfium2(pdf_path)

    assert "Could not extract text from PDF file" in str(exc_info.value)
    assert str(pdf_path) in str(exc_info.value.context["file_path"])
    assert "error" in exc_info.value.context
